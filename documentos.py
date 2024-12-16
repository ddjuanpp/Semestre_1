# documentos.py

from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

class DocumentUploader:
    def __init__(self):
        self.documents = []

    def add_document(self, file):
        """Leer documento y guardar su texto"""
        try:
            if file.name.endswith(".pdf"):
                text = self._extract_text_from_pdf(file)
            elif file.name.endswith(".docx"):
                text = self._extract_text_from_docx(file)
            elif file.name.endswith(".pptx"):
                text = self._extract_text_from_pptx(file)
            else:
                raise ValueError("Formato de archivo no soportado.")
            self.documents.append(text)
        except Exception as e:
            raise ValueError(f"Error al procesar el archivo {file.name}: {e}")

    def _extract_text_from_pdf(self, file):
        """Extraer texto de un archivo PDF"""
        from PyPDF2 import PdfReader
        reader = PdfReader(file)
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""  # Agrega el texto extraído de cada página
        return text

    def _extract_text_from_docx(self, file):
        """Extraer texto de un archivo DOCX"""
        from docx import Document
        doc = Document(file)
        return "\n".join([para.text for para in doc.paragraphs])

    def _extract_text_from_pptx(self, file):
        """Extraer texto de un archivo PPTX"""
        from pptx import Presentation
        presentation = Presentation(file)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text + "\n"
        return text

    def get_documents(self):
        """Retornar lista de documentos cargados"""
        return self.documents

    def get_concatenated_text(self):
        """Concatenar texto de todos los documentos"""
        return " ".join(self.documents)